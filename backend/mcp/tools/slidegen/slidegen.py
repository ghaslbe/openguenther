#!/usr/bin/env python3
"""
slidegen.py — AI-powered dark-themed presentation generator

Usage:
  python slidegen.py "Dein Thema"
  python slidegen.py "Machine Learning" --output ml.pptx
  python slidegen.py "Agile Methoden" --model anthropic/claude-3.5-haiku
  python slidegen.py "Klimawandel" --theme purple

Prerequisites:
  pip install python-pptx requests lxml
  export OPENROUTER_API_KEY="sk-or-..."

Layouts:
  hero        — Titel links, abstraktes Artwork rechts (Split)
  cards       — 3 Karten + optionale Fußnote
  two-column  — Bullet-Gruppen in 2 Spalten
  steps       — Nummerierte Prozessschritte mit Cards
  icon-list   — Icon-Kreise + Titel + Beschreibung
  pyramid     — Hierarchie-Pyramide
  feature     — Großer Titel + Icon-Liste links, Artwork rechts
  statement   — Artwork links, Titel + Content-Blöcke rechts
"""

import argparse
import json
import os
import re
import sys
from pathlib import Path

try:
    from lxml import etree
except ImportError:
    sys.exit("Error: 'lxml' not installed. Run: pip install lxml")

try:
    import requests
except ImportError:
    sys.exit("Error: 'requests' not installed. Run: pip install requests")

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.oxml.ns import qn
    from pptx.util import Inches, Pt
except ImportError:
    sys.exit("Error: 'python-pptx' not installed. Run: pip install python-pptx")


# ── Dimensions ────────────────────────────────────────────────────────────────
W   = Inches(13.333)
H   = Inches(7.5)
PAD = Inches(0.72)

# ── Themes ────────────────────────────────────────────────────────────────────
THEMES = {
    "dark": {
        "bg":       RGBColor(0x1A, 0x1A, 0x1A),
        "card":     RGBColor(0x26, 0x26, 0x26),
        "card2":    RGBColor(0x2E, 0x2E, 0x2E),
        "primary":  RGBColor(0xE8, 0x8C, 0x28),
        "primary2": RGBColor(0xC4, 0x6E, 0x10),
        "text":     RGBColor(0xED, 0xED, 0xED),
        "muted":    RGBColor(0x8A, 0x8A, 0x8A),
        "border":   RGBColor(0x3A, 0x3A, 0x3A),
        "deco":     RGBColor(0x24, 0x24, 0x24),
    },
    "purple": {
        "bg":       RGBColor(0x12, 0x0D, 0x24),
        "card":     RGBColor(0x1D, 0x16, 0x38),
        "card2":    RGBColor(0x27, 0x1E, 0x4A),
        "primary":  RGBColor(0xA0, 0x40, 0xD8),
        "primary2": RGBColor(0x7A, 0x28, 0xAC),
        "text":     RGBColor(0xEA, 0xE8, 0xF5),
        "muted":    RGBColor(0x7A, 0x72, 0x99),
        "border":   RGBColor(0x38, 0x2E, 0x58),
        "deco":     RGBColor(0x1A, 0x13, 0x30),
    },
}
T = THEMES["dark"]

FONT_HEAD = "Calibri"
FONT_BODY = "Calibri Light"


# ── XML / alpha helpers ───────────────────────────────────────────────────────

def _set_alpha(shape, pct: float):
    """Set solid-fill opacity. pct: 0=invisible, 100=opaque."""
    try:
        spPr = shape._element.spPr
        sf = spPr.find(qn("a:solidFill"))
        if sf is None:
            return
        clr = sf.find(qn("a:srgbClr"))
        if clr is None:
            return
        for old in clr.findall(qn("a:alpha")):
            clr.remove(old)
        el = etree.SubElement(clr, qn("a:alpha"))
        el.set("val", str(int(pct * 1000)))
    except Exception:
        pass


# ── Shape primitives ──────────────────────────────────────────────────────────

def rect(slide, l, t, w, h, fill, line=None, lw=Pt(0.5), alpha=100):
    sh = slide.shapes.add_shape(1, int(l), int(t), int(w), int(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line:
        sh.line.color.rgb = line
        sh.line.width = int(lw)
    else:
        sh.line.fill.background()
    if alpha < 100:
        _set_alpha(sh, alpha)
    return sh


def rrect(slide, l, t, w, h, fill, line=None, lw=Pt(0.5), radius=0.05, alpha=100):
    sh = slide.shapes.add_shape(5, int(l), int(t), int(w), int(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line:
        sh.line.color.rgb = line
        sh.line.width = int(lw)
    else:
        sh.line.fill.background()
    sh.adjustments[0] = radius
    if alpha < 100:
        _set_alpha(sh, alpha)
    return sh


def oval(slide, cx, cy, size, fill, alpha=100):
    s = int(size)
    sh = slide.shapes.add_shape(9, int(cx - s / 2), int(cy - s / 2), s, s)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    if alpha < 100:
        _set_alpha(sh, alpha)
    return sh


def txt(slide, content, l, t, w, h,
        size=13, bold=False, color=None, font=None,
        align=PP_ALIGN.LEFT, wrap=True, italic=False):
    tx = slide.shapes.add_textbox(int(l), int(t), int(w), int(h))
    tf = tx.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = content
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color if color else T["text"]
    run.font.name = font if font else FONT_BODY
    return tx


# ── Reusable components ───────────────────────────────────────────────────────

def dot_grid(slide, ol, ot, cols, rows, spacing, dot_r, color, alpha=20):
    for r in range(rows):
        for c in range(cols):
            oval(slide,
                 ol + c * spacing,
                 ot + r * spacing,
                 dot_r, color, alpha=alpha)


def abstract_panel(slide, l, t, w, h):
    """
    Decorative abstract artwork panel — replaces images.
    Inspired by the Gamma/modern-tool aesthetic (overlapping circles, dot grid).
    """
    # Panel background
    rrect(slide, l, t, w, h, T["card"], line=T["border"], lw=Pt(0.4), radius=0.03)

    cx, cy = l + w * 0.62, t + h * 0.38

    # Largest circle — very faint
    oval(slide, cx, cy,           min(w, h) * 1.05, T["primary"], alpha=9)
    # Mid circle — slightly more visible
    oval(slide, cx - w * 0.12, cy + h * 0.18, min(w, h) * 0.68, T["primary"], alpha=14)
    # Smaller accent circle top-left of panel
    oval(slide, l + w * 0.22, t + h * 0.24, min(w, h) * 0.38, T["primary"], alpha=20)
    # Tiny solid accent circle
    oval(slide, l + w * 0.72, t + h * 0.72, min(w, h) * 0.14, T["primary"], alpha=55)
    oval(slide, l + w * 0.30, t + h * 0.78, min(w, h) * 0.09, T["primary"], alpha=40)

    # Inner ring-like circle (card2 so it's just slightly lighter than bg)
    oval(slide, cx, cy, min(w, h) * 0.52, T["card2"], alpha=100)

    # Dot grid in the upper-left area of the panel
    dot_grid(slide,
             l + Inches(0.22), t + Inches(0.20),
             cols=4, rows=5,
             spacing=Inches(0.26),
             dot_r=Inches(0.048),
             color=T["primary"], alpha=22)

    # Thin geometric accent lines
    rect(slide, l + w * 0.14, t + h * 0.60, w * 0.55, Inches(0.018), T["border"], alpha=70)
    rect(slide, l + w * 0.24, t + h * 0.68, w * 0.38, Inches(0.018), T["border"], alpha=50)


def slide_header(slide, title, category=None, top=Inches(0.30),
                 left=None, right_edge=None):
    """
    Optional category pill + thin accent line + large bold title.
    Returns bottom y of the header area.
    `left` / `right_edge` let right-column layouts position the header freely.
    """
    l = left        if left        is not None else PAD
    r = right_edge  if right_edge  is not None else W - PAD
    w = r - l

    if category:
        pill_h = Inches(0.30)
        pill_w = max(Inches(1.0), len(category) * Inches(0.115) + Inches(0.28))
        rrect(slide, l, top, pill_w, pill_h, T["card2"], radius=0.40)
        oval(slide, l + Inches(0.14), top + pill_h / 2, Inches(0.09), T["primary"])
        txt(slide, category.upper(),
            l + Inches(0.27), top + Inches(0.03),
            pill_w - Inches(0.30), Inches(0.24),
            size=10, color=T["muted"], font=FONT_BODY)
        top += pill_h + Inches(0.08)

    rect(slide, l, top, min(Inches(1.8), w * 0.35), Inches(0.040), T["primary"])
    txt(slide, title,
        l, top + Inches(0.055),
        w, Inches(0.62),
        size=30, bold=True, color=T["text"], font=FONT_HEAD)
    return top + Inches(0.055) + Inches(0.64)


def _blank_layout(prs):
    for layout in prs.slide_layouts:
        if layout.name.lower() in ("blank", "leer", "vide", "vacío"):
            return layout
    return min(prs.slide_layouts, key=lambda l: len(list(l.placeholders)))


def _set_bg(slide):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = T["bg"]


def _rgb_parts(c: RGBColor):
    s = str(c)
    return int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16)


def _blend(c1: RGBColor, c2: RGBColor, t: float) -> RGBColor:
    r1, g1, b1 = _rgb_parts(c1)
    r2, g2, b2 = _rgb_parts(c2)
    return RGBColor(
        max(0, min(255, int(r1 + (r2 - r1) * t))),
        max(0, min(255, int(g1 + (g2 - g1) * t))),
        max(0, min(255, int(b1 + (b2 - b1) * t))),
    )

ICON_MAP = {
    "star": "★", "heart": "♥", "zap": "⚡", "shield": "◈",
    "target": "◎", "users": "⊕", "code": "{ }", "lightbulb": "✦",
    "handshake": "⊙", "graduation": "◉",
}


def render_lead(slide, text: str, left, top, width) -> float:
    """Render a lead/intro paragraph. Returns the Y position after the text."""
    if not text:
        return top
    h = Inches(0.60)
    rect(slide, left, top, Inches(0.06), h * 0.72, T["primary"], alpha=55)
    txt(slide, text, left + Inches(0.16), top, width - Inches(0.16), h,
        size=15, color=T["text"], font=FONT_BODY)
    return top + h + Inches(0.14)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE RENDERERS
# ══════════════════════════════════════════════════════════════════════════════

# ── HERO (split: text left | abstract art right) ──────────────────────────────
def render_hero(prs, data):
    """
    Left column: vertical accent bar + large title + subtitle + dots
    Right column: abstract art panel  (matches Screenshot 2 style)
    """
    slide = prs.slides.add_slide(_blank_layout(prs))
    _set_bg(slide)

    title    = data.get("title", "")
    subtitle = data.get("subtitle", "")

    split = W * 0.54   # left column width

    # ── Right: abstract art panel ────────────────────────────────────────────
    art_l = split + Inches(0.10)
    art_t = Inches(0.30)
    art_w = W - art_l - Inches(0.30)
    art_h = H - Inches(0.60)
    abstract_panel(slide, art_l, art_t, art_w, art_h)

    # ── Left: content ────────────────────────────────────────────────────────
    # Vertical accent bar
    bar_h = Inches(3.0)
    bar_y = (H - bar_h) / 2 - Inches(0.1)
    rrect(slide, PAD, bar_y, Inches(0.20), bar_h, T["primary"], radius=0.04)

    # Title (large, multi-line capable)
    txt(slide, title,
        PAD + Inches(0.42), H * 0.18,
        split - PAD - Inches(0.55), Inches(3.0),
        size=46, bold=True, color=T["text"], font=FONT_HEAD)

    # Subtitle
    txt(slide, subtitle,
        PAD + Inches(0.42), H * 0.62,
        split - PAD - Inches(0.55), Inches(1.1),
        size=17, color=T["muted"], font=FONT_BODY)

    # Tag row (from KI-Generierung PPTX slide 1: pill tags below subtitle)
    tags = data.get("tags", [])
    if tags:
        tx = PAD + Inches(0.42)
        ty = H * 0.62 + Inches(1.15)
        for tag in tags[:4]:
            pill_w = max(Inches(0.9), len(tag) * Inches(0.11) + Inches(0.30))
            rrect(slide, tx, ty, pill_w, Inches(0.30), T["card2"], radius=0.40)
            oval(slide, tx + Inches(0.14), ty + Inches(0.15),
                 Inches(0.08), T["primary"])
            txt(slide, tag.upper(),
                tx + Inches(0.26), ty + Inches(0.03),
                pill_w - Inches(0.28), Inches(0.24),
                size=10, color=T["muted"], font=FONT_BODY)
            tx += pill_w + Inches(0.12)

    # Bottom strip
    rect(slide, 0, H - Inches(0.08), W, Inches(0.08), T["primary"])


# ── CARDS (clean: no badges, optional footnote) ───────────────────────────────
def render_cards(prs, data):
    """
    Matches Screenshot 1: clean cards with bold title, muted description,
    optional footnote below the card row.
    """
    slide = prs.slides.add_slide(_blank_layout(prs))
    _set_bg(slide)

    title    = data.get("title", "")
    lead     = data.get("lead", "")
    footnote = data.get("subtitle", data.get("footnote", ""))
    cards    = data.get("cards", [])[:3]
    n = len(cards)
    if n == 0:
        return

    # Subtle deco circle top-right
    oval(slide, W * 0.93, H * 0.08, Inches(3.2), T["primary"], alpha=6)

    ct = slide_header(slide, title)
    ct = render_lead(slide, lead, PAD, ct, W - PAD * 2)

    footnote_h = Inches(0.55) if footnote else 0
    gap = Inches(0.22)
    cw  = (W - PAD * 2 - gap * (n - 1)) / n
    ch  = H - ct - Inches(0.30) - footnote_h

    for i, card in enumerate(cards):
        cx = PAD + i * (cw + gap)
        cy = ct

        # Shadow
        rrect(slide,
              cx + Inches(0.05), cy + Inches(0.07),
              cw, ch,
              RGBColor(0x0C, 0x0C, 0x0C), radius=0.05, alpha=65)

        # Card — matches frontend: bg-accent border border-border rounded-lg
        rrect(slide, cx, cy, cw, ch,
              T["card"], line=T["border"], lw=Pt(0.5), radius=0.05)

        # Thin left accent stripe
        rrect(slide, cx, cy, Inches(0.07), ch, T["primary"], radius=0.05)
        rect(slide,  cx, cy, Inches(0.07), ch * 0.3, T["primary"])  # fill top-rounded corner

        # Card title — bold, white (matches: text-xl font-semibold text-foreground)
        txt(slide, card.get("title", ""),
            cx + Inches(0.22), cy + Inches(0.20),
            cw - Inches(0.36), Inches(0.68),
            size=21, bold=True, color=T["text"], font=FONT_HEAD)

        # Thin separator
        rect(slide,
             cx + Inches(0.22), cy + Inches(0.90),
             cw - Inches(0.38), Inches(0.025), T["border"])

        # Description — muted (matches: text-sm text-muted-foreground)
        txt(slide, card.get("description", ""),
            cx + Inches(0.22), cy + Inches(0.98),
            cw - Inches(0.36), ch - Inches(1.08),
            size=14, color=T["muted"], font=FONT_BODY)

    # Optional footnote below cards (matches: mt-8 text-muted-foreground text-sm)
    if footnote:
        fy = H - footnote_h - Inches(0.15)
        rect(slide, PAD, fy - Inches(0.04), Inches(0.8), Inches(0.03), T["primary"], alpha=60)
        txt(slide, footnote,
            PAD, fy,
            W - PAD * 2, Inches(0.50),
            size=12.5, color=T["muted"], font=FONT_BODY, italic=True)


# ── TWO-COLUMN (bullet groups in 2 columns) ───────────────────────────────────
def render_two_column(prs, data):
    slide = prs.slides.add_slide(_blank_layout(prs))
    _set_bg(slide)

    ct = slide_header(slide, data.get("title", ""))
    ct = render_lead(slide, data.get("lead", ""), PAD, ct, W - PAD * 2)
    oval(slide, W * 0.92, H * 0.88, Inches(2.8), T["primary"], alpha=6)

    groups  = data.get("bullet_groups", data.get("groups", []))
    mid     = (len(groups) + 1) // 2
    cols    = [groups[:mid], groups[mid:]]
    col_gap = Inches(0.55)
    col_w   = (W - PAD * 2 - col_gap) / 2

    # Vertical divider
    div_x = PAD + col_w + col_gap / 2 - Inches(0.011)
    rect(slide, div_x, ct, Inches(0.022), H - ct - Inches(0.60), T["border"])

    for col_idx, col_groups in enumerate(cols):
        cx = PAD + col_idx * (col_w + col_gap)
        cy = ct

        for group in col_groups:
            if cy >= H - Inches(0.9):
                break

            heading = group.get("heading", group.get("title", ""))
            bullets = group.get("bullets", [])

            # Heading with small rounded left accent
            rrect(slide, cx, cy + Inches(0.06), Inches(0.09), Inches(0.30),
                  T["primary"], radius=0.20)
            txt(slide, heading,
                cx + Inches(0.17), cy,
                col_w - Inches(0.20), Inches(0.46),
                size=17, bold=True, color=T["primary"], font=FONT_HEAD)
            cy += Inches(0.50)

            for bullet in bullets:
                if cy >= H - Inches(0.75):
                    break
                oval(slide,
                     cx + Inches(0.13), cy + Inches(0.175),
                     Inches(0.092), T["primary"])
                txt(slide, bullet,
                    cx + Inches(0.30), cy,
                    col_w - Inches(0.36), Inches(0.46),
                    size=14, color=T["muted"], font=FONT_BODY)
                cy += Inches(0.44)

            cy += Inches(0.28)


# ── STEPS (numbered cards) ────────────────────────────────────────────────────
def render_steps(prs, data):
    slide = prs.slides.add_slide(_blank_layout(prs))
    _set_bg(slide)

    ct = slide_header(slide, data.get("title", ""))
    ct = render_lead(slide, data.get("lead", ""), PAD, ct, W - PAD * 2)
    oval(slide, W * 0.91, H * 0.5, Inches(3.5), T["primary"], alpha=5)

    steps   = data.get("steps", [])[:5]
    n       = len(steps)
    if n == 0:
        return

    avail_h = H - ct - Inches(0.55)
    step_h  = avail_h / n
    badge_s = Inches(0.48)
    badge_x = PAD + badge_s / 2

    for i, step in enumerate(steps):
        sy  = ct + i * step_h
        bcy = sy + step_h / 2

        # Connector line
        if i < n - 1:
            line_top = bcy + badge_s / 2
            line_bot = bcy + step_h
            rect(slide, badge_x - Inches(0.020), line_top,
                 Inches(0.040), line_bot - line_top,
                 T["primary"], alpha=30)

        # Step card
        card_t = sy + Inches(0.07)
        card_h = step_h - Inches(0.12)
        rrect(slide, PAD, card_t, W - PAD * 2, card_h,
              T["card"], line=T["border"], lw=Pt(0.4), radius=0.04)

        # Left accent bar on card
        rrect(slide, PAD, card_t, Inches(0.07), card_h,
              T["primary"], radius=0.04)
        rect(slide, PAD, card_t, Inches(0.07), card_h * 0.25, T["primary"])

        # Circle badge
        oval(slide, badge_x, bcy, badge_s, T["primary"])
        txt(slide, str(i + 1),
            badge_x - badge_s / 2, bcy - badge_s / 2,
            badge_s, badge_s,
            size=15, bold=True, color=T["bg"], font=FONT_HEAD,
            align=PP_ALIGN.CENTER)

        # Text
        tx  = PAD + Inches(0.20) + badge_s
        tw  = W - tx - PAD - Inches(0.25)
        ty0 = card_t + (card_h - Inches(0.82)) / 2

        txt(slide, step.get("title", ""),
            tx, ty0, tw, Inches(0.46),
            size=18, bold=True, color=T["text"], font=FONT_HEAD)
        txt(slide, step.get("description", ""),
            tx, ty0 + Inches(0.46), tw, Inches(0.42),
            size=14, color=T["muted"], font=FONT_BODY)


# ── ICON LIST (2-column grid with circle icons) ───────────────────────────────
def render_icon_list(prs, data):
    slide = prs.slides.add_slide(_blank_layout(prs))
    _set_bg(slide)

    ct = slide_header(slide, data.get("title", ""))

    items   = data.get("items", [])[:6]
    n       = len(items)
    if n == 0:
        return

    cols    = 2
    rows    = (n + 1) // 2
    col_gap = Inches(0.38)
    col_w   = (W - PAD * 2 - col_gap) / 2
    avail_h = H - ct - Inches(0.50)
    row_h   = avail_h / rows
    icon_s  = Inches(0.52)
    ipad    = Inches(0.22)

    for i, item in enumerate(items):
        col = i % cols
        row = i // cols
        ix  = PAD + col * (col_w + col_gap)
        iy  = ct + row * row_h
        ch  = row_h - Inches(0.13)

        # Item card
        rrect(slide, ix, iy, col_w, ch,
              T["card"], line=T["border"], lw=Pt(0.4), radius=0.06)

        # Double-circle icon
        icx = ix + ipad + icon_s / 2
        icy = iy + ch / 2
        oval(slide, icx, icy, icon_s,          T["primary"])
        oval(slide, icx, icy, icon_s * 0.62,   T["primary2"])

        symbol = ICON_MAP.get(item.get("icon", ""), "●")
        txt(slide, symbol,
            icx - icon_s / 2, icy - icon_s / 2,
            icon_s, icon_s,
            size=16, bold=True, color=T["text"], font=FONT_BODY,
            align=PP_ALIGN.CENTER)

        tx  = ix + ipad + icon_s + Inches(0.18)
        tw  = col_w - ipad - icon_s - Inches(0.30)
        ty0 = iy + (ch - Inches(0.78)) / 2

        txt(slide, item.get("title", ""),
            tx, ty0, tw, Inches(0.40),
            size=17, bold=True, color=T["text"], font=FONT_HEAD)
        txt(slide, item.get("description", ""),
            tx, ty0 + Inches(0.40), tw, Inches(0.44),
            size=13, color=T["muted"], font=FONT_BODY)

        oval(slide,
             ix + col_w - Inches(0.26), iy + ch - Inches(0.26),
             Inches(0.09), T["primary"], alpha=38)


# ── PYRAMID ───────────────────────────────────────────────────────────────────
def render_pyramid(prs, data):
    slide = prs.slides.add_slide(_blank_layout(prs))
    _set_bg(slide)

    ct = slide_header(slide, data.get("title", ""))
    oval(slide, PAD * 0.5, H * 0.5, Inches(2.5), T["primary"], alpha=5)

    items   = data.get("items", [])[:5]
    n       = len(items)
    if n == 0:
        return

    pyr_w   = W * 0.52
    desc_x  = PAD + pyr_w + Inches(0.38)
    desc_w  = W - desc_x - PAD
    avail_h = H - ct - Inches(0.50)
    tier_gap= Inches(0.07)
    tier_h  = (avail_h - tier_gap * (n - 1)) / n
    center_x= PAD + pyr_w / 2

    for i, item in enumerate(items):
        frac = (i + 1) / n
        tw   = pyr_w * (0.16 + 0.84 * frac)
        tx   = center_x - tw / 2
        ty   = ct + i * (tier_h + tier_gap)

        color = _blend(T["primary"], _blend(T["muted"], T["card2"], 0.5),
                       i / max(n - 1, 1) * 0.65)

        rrect(slide, tx, ty, tw, tier_h, color, radius=0.25)
        txt(slide, item.get("title", ""),
            tx + Inches(0.10), ty, tw - Inches(0.20), tier_h,
            size=12.5, bold=True, color=T["bg"], font=FONT_HEAD,
            align=PP_ALIGN.CENTER)

        # Connecting dot
        oval(slide, desc_x - Inches(0.26), ty + tier_h / 2,
             Inches(0.09), T["primary"], alpha=55)
        txt(slide, item.get("description", ""),
            desc_x, ty + (tier_h - Inches(0.36)) / 2,
            desc_w, Inches(0.38),
            size=12.5, color=T["muted"], font=FONT_BODY)

    # Vertical connecting line
    lt = ct + tier_h / 2
    lb = ct + (n - 1) * (tier_h + tier_gap) + tier_h / 2
    rect(slide, desc_x - Inches(0.261), lt, Inches(0.022), lb - lt, T["border"])


# ── FEATURE (Screenshot 3: title + icon-list left | art panel right) ──────────
def render_feature(prs, data):
    """
    Left half: big section title + subtitle + icon-list items with chevron style
    Right half: abstract art panel
    Matches the purple-theme example (Screenshot 3).
    """
    slide = prs.slides.add_slide(_blank_layout(prs))
    _set_bg(slide)

    title    = data.get("title", "")
    subtitle = data.get("subtitle", "")
    items    = data.get("items", [])[:4]

    split  = W * 0.50
    art_l  = split + Inches(0.15)
    art_t  = Inches(0.28)
    art_w  = W - art_l - Inches(0.28)
    art_h  = H - Inches(0.56)

    # Right: abstract art
    abstract_panel(slide, art_l, art_t, art_w, art_h)

    # Left: vertical accent bar
    bar_h = H - Inches(0.80)
    bar_y = Inches(0.40)
    rrect(slide, PAD, bar_y, Inches(0.18), bar_h, T["primary"], radius=0.04)

    # Title (large, multi-line)
    txt(slide, title,
        PAD + Inches(0.36), Inches(0.36),
        split - PAD - Inches(0.50), Inches(2.20),
        size=34, bold=True, color=T["text"], font=FONT_HEAD)

    # Subtitle
    if subtitle:
        txt(slide, subtitle,
            PAD + Inches(0.36), Inches(2.64),
            split - PAD - Inches(0.50), Inches(0.55),
            size=14, color=T["muted"], font=FONT_BODY)

    # Icon items
    item_top = Inches(3.30) if subtitle else Inches(2.80)
    avail    = H - item_top - Inches(0.45)
    item_h   = avail / max(len(items), 1)

    for i, item in enumerate(items):
        iy      = item_top + i * item_h
        icon_s  = Inches(0.38)
        icx     = PAD + Inches(0.36) + icon_s / 2
        icy     = iy + item_h / 2

        # Small icon circle
        oval(slide, icx, icy, icon_s, T["primary"])
        oval(slide, icx, icy, icon_s * 0.58, T["primary2"])
        symbol = ICON_MAP.get(item.get("icon", ""), "◈")
        txt(slide, symbol,
            icx - icon_s / 2, icy - icon_s / 2,
            icon_s, icon_s,
            size=12, bold=True, color=T["text"], font=FONT_BODY,
            align=PP_ALIGN.CENTER)

        # Title + description
        tx = PAD + Inches(0.36) + icon_s + Inches(0.16)
        tw = split - tx - Inches(0.30)
        txt(slide, item.get("title", ""),
            tx, iy + (item_h - Inches(0.76)) / 2,
            tw, Inches(0.38),
            size=16, bold=True, color=T["text"], font=FONT_HEAD)
        txt(slide, item.get("description", ""),
            tx, iy + (item_h - Inches(0.76)) / 2 + Inches(0.38),
            tw, Inches(0.40),
            size=13.5, color=T["muted"], font=FONT_BODY)

    # Bottom strip
    rect(slide, 0, H - Inches(0.07), W, Inches(0.07), T["primary"])


# ── STATEMENT (Screenshot 4: art left | bold title + content blocks right) ────
def render_statement(prs, data):
    """
    Left half: abstract art panel
    Right half: large title + content blocks (heading + body text each)
    Matches the purple-theme 'Branchen-Know-how' example (Screenshot 4).
    """
    slide = prs.slides.add_slide(_blank_layout(prs))
    _set_bg(slide)

    title  = data.get("title", "")
    blocks = data.get("blocks", [])[:4]

    split = W * 0.44
    art_t = Inches(0.28)
    art_h = H - Inches(0.56)

    # Left: abstract art
    abstract_panel(slide, Inches(0.28), art_t, split - Inches(0.28), art_h)

    # Right: content
    rx   = split + Inches(0.30)
    rw   = W - rx - PAD
    ry   = Inches(0.38)

    # Thin top accent line
    rect(slide, rx, ry, Inches(1.6), Inches(0.038), T["primary"])

    # Large title
    txt(slide, title,
        rx, ry + Inches(0.06),
        rw, Inches(2.10),
        size=32, bold=True, color=T["text"], font=FONT_HEAD)

    # Content blocks
    block_top = ry + Inches(2.30)
    avail_h   = H - block_top - Inches(0.50)
    block_h   = avail_h / max(len(blocks), 1)

    for i, block in enumerate(blocks):
        by = block_top + i * block_h

        # Small diamond/dot marker
        oval(slide, rx + Inches(0.10), by + Inches(0.16),
             Inches(0.10), T["primary"])

        txt(slide, block.get("heading", ""),
            rx + Inches(0.26), by,
            rw - Inches(0.26), Inches(0.42),
            size=16, bold=True, color=T["text"], font=FONT_HEAD)
        txt(slide, block.get("text", ""),
            rx + Inches(0.26), by + Inches(0.40),
            rw - Inches(0.26), block_h - Inches(0.44),
            size=13.5, color=T["muted"], font=FONT_BODY)


# ── STATS (Slide 2/5 style: big numbers right, content left) ─────────────────
def render_stats(prs, data):
    """
    Left column: category + title + body text + optional quote block.
    Right column: 2–4 big statistics, each with huge value, thin accent line,
                  label, and short description.
    Inspired by slides 2 & 5 of KI-fur-Code-Generierung2.pptx.
    """
    slide = prs.slides.add_slide(_blank_layout(prs))
    _set_bg(slide)

    category = data.get("category")
    title    = data.get("title", "")
    body     = data.get("body", "")
    stats    = data.get("stats", [])[:4]
    quote    = data.get("quote", "")
    author   = data.get("author", "")

    split = W * 0.50
    ct    = slide_header(slide, title, category=category)
    ct    = render_lead(slide, data.get("lead", ""), PAD, ct, split - PAD - Inches(0.35))

    # Deco circle (bottom-left)
    oval(slide, PAD * 0.6, H * 0.85, Inches(2.2), T["primary"], alpha=6)

    # Body text
    if body:
        txt(slide, body,
            PAD, ct,
            split - PAD - Inches(0.35), Inches(1.5),
            size=15, color=T["muted"], font=FONT_BODY)

    # Quote / callout box (bottom-left)
    if quote:
        qh  = Inches(1.20) if author else Inches(0.95)
        qy  = H - qh - Inches(0.35)
        rrect(slide, PAD, qy, split - PAD - Inches(0.35), qh,
              T["card"], line=T["border"], lw=Pt(0.4), radius=0.04)
        rrect(slide, PAD, qy, Inches(0.07), qh, T["primary"], radius=0.04)
        rect(slide,  PAD, qy, Inches(0.07), qh * 0.3, T["primary"])
        txt(slide, f'"{quote}"',
            PAD + Inches(0.18), qy + Inches(0.12),
            split - PAD - Inches(0.55), Inches(0.68),
            size=12, color=T["text"], font=FONT_BODY, italic=True)
        if author:
            txt(slide, f"— {author}",
                PAD + Inches(0.18), qy + qh - Inches(0.32),
                split - PAD - Inches(0.55), Inches(0.26),
                size=11, bold=True, color=T["muted"], font=FONT_BODY)

    # Right: big stats
    rx      = split + Inches(0.20)
    rw      = W - rx - Inches(0.40)
    n       = len(stats)
    avail_h = H - ct - Inches(0.40)
    item_h  = avail_h / max(n, 1)

    for i, stat in enumerate(stats):
        sy = ct + i * item_h

        # Big value
        txt(slide, stat.get("value", ""),
            rx, sy,
            rw, Inches(0.85),
            size=54, bold=True, color=T["primary"], font=FONT_HEAD)

        # Thin accent rule below value
        rect(slide, rx, sy + Inches(0.85), rw * 0.50, Inches(0.038), T["primary"])

        # Label (bold, white)
        txt(slide, stat.get("label", ""),
            rx, sy + Inches(0.92),
            rw, Inches(0.40),
            size=17, bold=True, color=T["text"], font=FONT_HEAD)

        # Description (muted)
        txt(slide, stat.get("description", ""),
            rx, sy + Inches(1.34),
            rw, item_h - Inches(1.40),
            size=13.5, color=T["muted"], font=FONT_BODY)


# ── TOOLLIST (Slide 6 style: art-panel left, stacked tool items right) ─────────
def render_toollist(prs, data):
    """
    Left: abstract art panel.
    Right: category + title + 3–4 vertically stacked items, each with
           left accent stripe, title, and description.
    Inspired by slide 6 of KI-fur-Code-Generierung2.pptx.
    """
    slide = prs.slides.add_slide(_blank_layout(prs))
    _set_bg(slide)

    category = data.get("category")
    title    = data.get("title", "")
    items    = data.get("items", [])[:4]

    split = W * 0.31
    art_t = Inches(0.28)
    art_h = H - Inches(0.56)
    abstract_panel(slide, Inches(0.26), art_t, split - Inches(0.26), art_h)

    rx  = split + Inches(0.22)
    rw  = W - rx - PAD * 0.6
    ct  = slide_header(slide, title, category=category,
                       left=rx, right_edge=W - Inches(0.35), top=Inches(0.28))

    n       = len(items)
    avail_h = H - ct - Inches(0.40)
    gap     = Inches(0.16)
    item_h  = (avail_h - gap * (n - 1)) / max(n, 1)

    for i, item in enumerate(items):
        ix = rx
        iy = ct + i * (item_h + gap)

        # Shadow
        rrect(slide, ix + Inches(0.04), iy + Inches(0.05),
              rw, item_h, RGBColor(0x0A, 0x0A, 0x0A), radius=0.04, alpha=60)
        # Card
        rrect(slide, ix, iy, rw, item_h,
              T["card"], line=T["border"], lw=Pt(0.4), radius=0.04)
        # Left accent stripe (thick, like the PPTX)
        rrect(slide, ix, iy, Inches(0.12), item_h, T["primary"], radius=0.04)
        rect(slide,  ix, iy, Inches(0.12), item_h * 0.28, T["primary"])

        # Title
        txt(slide, item.get("title", ""),
            ix + Inches(0.24), iy + Inches(0.12),
            rw - Inches(0.36), Inches(0.40),
            size=17, bold=True, color=T["text"], font=FONT_HEAD)

        # Separator
        rect(slide, ix + Inches(0.24), iy + Inches(0.54),
             rw - Inches(0.40), Inches(0.020), T["border"])

        # Description
        txt(slide, item.get("description", ""),
            ix + Inches(0.24), iy + Inches(0.60),
            rw - Inches(0.36), item_h - Inches(0.68),
            size=14, color=T["muted"], font=FONT_BODY)


# ── ROADMAP (Slide 7 style: two columns of numbered items + CTA banner) ────────
def render_roadmap(prs, data):
    """
    Two columns, each with a heading and 2–4 numbered items.
    Each item: small number badge → thin accent rule → title → description.
    Optional CTA banner at the bottom.
    Inspired by slide 7 of KI-fur-Code-Generierung2.pptx.
    """
    slide = prs.slides.add_slide(_blank_layout(prs))
    _set_bg(slide)

    category = data.get("category")
    title    = data.get("title", "")
    columns  = data.get("columns", [])[:2]
    cta      = data.get("cta", "")

    cta_h = Inches(0.98) if cta else 0
    ct    = slide_header(slide, title, category=category)
    ct    = render_lead(slide, data.get("lead", ""), PAD, ct, W - PAD * 2)

    # Subtle deco
    oval(slide, W * 0.88, H * 0.15, Inches(2.8), T["primary"], alpha=5)

    col_gap = Inches(0.42)
    col_w   = (W - PAD * 2 - col_gap) / 2
    bot     = H - cta_h - Inches(0.30)

    # Vertical divider
    div_x = PAD + col_w + col_gap / 2 - Inches(0.011)
    rect(slide, div_x, ct, Inches(0.022), bot - ct, T["border"])

    for col_idx, col in enumerate(columns):
        cx     = PAD + col_idx * (col_w + col_gap)
        cy     = ct
        items  = col.get("items", [])[:4]
        accent = T["primary"] if col_idx == 0 else T["muted"]

        # Column heading
        txt(slide, col.get("heading", ""),
            cx, cy, col_w, Inches(0.42),
            size=17, bold=True, color=T["text"], font=FONT_HEAD)
        cy += Inches(0.48)

        for item in items:
            if cy >= bot - Inches(0.55):
                break

            num    = item.get("number", str(items.index(item) + 1).zfill(2))
            num_w  = Inches(0.34)

            # Number badge
            rrect(slide, cx, cy, num_w, Inches(0.25),
                  T["card2"], radius=0.25)
            txt(slide, num,
                cx, cy, num_w, Inches(0.25),
                size=10, bold=True, color=accent, font=FONT_HEAD,
                align=PP_ALIGN.CENTER)

            # Rule line (primary for left, muted for right)
            rule_x = cx + num_w + Inches(0.10)
            rule_w = col_w - num_w - Inches(0.10)
            rect(slide, rule_x, cy + Inches(0.11),
                 rule_w, Inches(0.025), accent)
            cy += Inches(0.32)

            # Title
            txt(slide, item.get("title", ""),
                cx, cy, col_w, Inches(0.40),
                size=16, bold=True, color=T["text"], font=FONT_HEAD)
            cy += Inches(0.40)

            # Description
            desc_h = Inches(0.50)
            txt(slide, item.get("description", ""),
                cx, cy, col_w, desc_h,
                size=13.5, color=T["muted"], font=FONT_BODY)
            cy += desc_h + Inches(0.20)

    # CTA Banner
    if cta:
        cta_y = H - cta_h - Inches(0.18)
        rrect(slide, PAD, cta_y, W - PAD * 2, cta_h,
              T["primary"], alpha=18, radius=0.03)
        rrect(slide, PAD, cta_y, Inches(0.08), cta_h,
              T["primary"], radius=0.03)
        rect(slide,  PAD, cta_y, Inches(0.08), cta_h * 0.30, T["primary"])
        txt(slide, cta,
            PAD + Inches(0.22), cta_y + (cta_h - Inches(0.44)) / 2,
            W - PAD * 2 - Inches(0.32), Inches(0.46),
            size=16, bold=True, color=T["text"], font=FONT_HEAD)


# ── COMPARISON (Slide 8 style: two contrasting boxes) ─────────────────────────
def render_comparison(prs, data):
    """
    Two large boxes side by side:
      Left  — positive / pros  (primary accent)
      Right — challenges / cons (muted)
    Inspired by the Vorteile vs. Herausforderungen slide in the example PPTX.
    """
    slide = prs.slides.add_slide(_blank_layout(prs))
    _set_bg(slide)

    category = data.get("category")
    ct = slide_header(slide, data.get("title", ""), category=category)

    left  = data.get("left",  {})
    right = data.get("right", {})

    gap = Inches(0.28)
    bw  = (W - PAD * 2 - gap) / 2
    bh  = H - ct - Inches(0.45)

    # ── Left box (positive) ───────────────────────────────────────────────────
    lx = PAD
    # Shadow
    rrect(slide, lx + Inches(0.05), ct + Inches(0.07),
          bw, bh, RGBColor(0x0A, 0x0A, 0x0A), radius=0.04, alpha=65)
    # Box — slightly lighter card, primary border
    rrect(slide, lx, ct, bw, bh,
          T["card"], line=T["primary"], lw=Pt(1.0), radius=0.04)
    # Left accent stripe
    rrect(slide, lx, ct, Inches(0.09), bh, T["primary"], radius=0.04)
    rect(slide,  lx, ct, Inches(0.09), bh * 0.28, T["primary"])

    # Header row
    hdr_h = Inches(0.56)
    rrect(slide, lx, ct, bw, hdr_h, T["primary"], alpha=18, radius=0.04)
    rect(slide,  lx, ct + hdr_h * 0.5, bw, hdr_h * 0.5, T["primary"], alpha=18)

    # Check-circle icon
    oval(slide, lx + Inches(0.34), ct + hdr_h / 2, Inches(0.34), T["primary"])
    txt(slide, "✓",
        lx + Inches(0.18), ct + Inches(0.08), Inches(0.34), Inches(0.34),
        size=13, bold=True, color=T["bg"], font=FONT_HEAD, align=PP_ALIGN.CENTER)

    # Label
    txt(slide, left.get("label", "Vorteile").upper(),
        lx + Inches(0.60), ct + Inches(0.10),
        bw - Inches(0.72), Inches(0.36),
        size=14, bold=True, color=T["primary"], font=FONT_HEAD)

    # Items
    iy = ct + hdr_h + Inches(0.18)
    for item in left.get("items", [])[:6]:
        oval(slide, lx + Inches(0.27), iy + Inches(0.165),
             Inches(0.10), T["primary"])
        txt(slide, item,
            lx + Inches(0.46), iy,
            bw - Inches(0.60), Inches(0.42),
            size=14, color=T["text"], font=FONT_BODY)
        iy += Inches(0.46)

    # ── Right box (challenges) ─────────────────────────────────────────────────
    rx = PAD + bw + gap
    # Shadow
    rrect(slide, rx + Inches(0.05), ct + Inches(0.07),
          bw, bh, RGBColor(0x0A, 0x0A, 0x0A), radius=0.04, alpha=65)
    # Box — slightly different shade, muted border
    rrect(slide, rx, ct, bw, bh,
          T["card2"], line=T["border"], lw=Pt(0.5), radius=0.04)

    # Header row (muted)
    rrect(slide, rx, ct, bw, hdr_h, T["border"], alpha=40, radius=0.04)
    rect(slide,  rx, ct + hdr_h * 0.5, bw, hdr_h * 0.5, T["border"], alpha=40)

    # Warning circle icon
    oval(slide, rx + Inches(0.34), ct + hdr_h / 2, Inches(0.34), T["muted"], alpha=60)
    txt(slide, "!",
        rx + Inches(0.18), ct + Inches(0.08), Inches(0.34), Inches(0.34),
        size=13, bold=True, color=T["muted"], font=FONT_HEAD, align=PP_ALIGN.CENTER)

    # Label
    txt(slide, right.get("label", "Herausforderungen").upper(),
        rx + Inches(0.60), ct + Inches(0.10),
        bw - Inches(0.72), Inches(0.36),
        size=14, bold=True, color=T["muted"], font=FONT_HEAD)

    # Items
    iy = ct + hdr_h + Inches(0.18)
    for item in right.get("items", [])[:6]:
        rect(slide, rx + Inches(0.22), iy + Inches(0.19),
             Inches(0.10), Inches(0.03), T["muted"], alpha=70)
        txt(slide, item,
            rx + Inches(0.40), iy,
            bw - Inches(0.54), Inches(0.42),
            size=14, color=T["muted"], font=FONT_BODY)
        iy += Inches(0.46)


# ── GRID (Slide 4/10 style: 2×2 cards with left accent stripe) ────────────────
def render_grid(prs, data):
    """
    2×2 card grid. Each card has a thin left accent stripe, bold title,
    separator line, and description.
    Inspired by slides 4 & 10 in the example PPTX.
    """
    slide = prs.slides.add_slide(_blank_layout(prs))
    _set_bg(slide)

    category = data.get("category")
    ct = slide_header(slide, data.get("title", ""), category=category)

    # Subtle deco
    oval(slide, W * 0.92, H * 0.10, Inches(3.0), T["primary"], alpha=6)

    cards = data.get("cards", [])[:4]
    n     = len(cards)
    if n == 0:
        return

    cols    = 2
    rows    = (n + cols - 1) // cols
    h_gap   = Inches(0.22)
    v_gap   = Inches(0.20)
    cw      = (W - PAD * 2 - h_gap) / 2
    avail_h = H - ct - Inches(0.40)
    ch      = (avail_h - v_gap * (rows - 1)) / rows

    for i, card in enumerate(cards):
        col = i % cols
        row = i // cols
        cx  = PAD + col * (cw + h_gap)
        cy  = ct + row * (ch + v_gap)

        # Shadow
        rrect(slide, cx + Inches(0.05), cy + Inches(0.06),
              cw, ch, RGBColor(0x0A, 0x0A, 0x0A), radius=0.04, alpha=60)
        # Card body
        rrect(slide, cx, cy, cw, ch,
              T["card"], line=T["border"], lw=Pt(0.4), radius=0.04)

        # Left accent stripe — exact style from the PPTX (thin bar, full height)
        rrect(slide, cx, cy, Inches(0.10), ch, T["primary"], radius=0.04)
        rect(slide,  cx, cy, Inches(0.10), ch * 0.25, T["primary"])

        # Title
        txt(slide, card.get("title", ""),
            cx + Inches(0.22), cy + Inches(0.18),
            cw - Inches(0.36), Inches(0.50),
            size=18, bold=True, color=T["text"], font=FONT_HEAD)

        # Separator
        rect(slide,
             cx + Inches(0.22), cy + Inches(0.70),
             cw - Inches(0.38), Inches(0.022), T["border"])

        # Description
        txt(slide, card.get("description", ""),
            cx + Inches(0.22), cy + Inches(0.76),
            cw - Inches(0.36), ch - Inches(0.86),
            size=13, color=T["muted"], font=FONT_BODY)

        # Small accent dot bottom-right
        oval(slide,
             cx + cw - Inches(0.28), cy + ch - Inches(0.28),
             Inches(0.10), T["primary"], alpha=35)


# ── Renderer dispatch ─────────────────────────────────────────────────────────

RENDERERS = {
    "hero":        render_hero,
    "cards":       render_cards,
    "two-column":  render_two_column,
    "two_column":  render_two_column,
    "steps":       render_steps,
    "icon-list":   render_icon_list,
    "icon_list":   render_icon_list,
    "pyramid":     render_pyramid,
    "feature":     render_feature,
    "statement":   render_statement,
    "comparison":  render_comparison,
    "grid":        render_grid,
    "stats":       render_stats,
    "toollist":    render_toollist,
    "roadmap":     render_roadmap,
}


def build_pptx(slide_list: list) -> Presentation:
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    for slide_data in slide_list:
        stype    = slide_data.get("type", "")
        renderer = RENDERERS.get(stype)
        label    = slide_data.get("title", stype)
        if renderer:
            renderer(prs, slide_data)
            print(f"  + [{stype:<12}] {label}")
        else:
            print(f"  ! Unknown type '{stype}' — skipped")

    return prs


# ── LLM ──────────────────────────────────────────────────────────────────────

SYSTEM_PROMPT = """\
You are a professional presentation designer. Generate a complete presentation as valid JSON.

Available slide types and their required fields:

  "hero"      → {type, title, subtitle}
               Split layout: title left, abstract artwork right.

  "cards"     → {type, title, lead?, cards:[{title,description}×3], subtitle?}
               3 clean cards + optional footnote/subtitle below them.
               lead = 1-sentence intro shown above the cards.

  "two-column"→ {type, title, lead?, bullet_groups:[{heading, bullets:[str]}]}
               Bullet-point groups arranged in 2 columns.
               lead = 1-sentence intro shown above the columns.

  "steps"     → {type, title, lead?, steps:[{title,description}]}
               Numbered process steps with cards.
               lead = 1-sentence intro shown above the steps.

  "icon-list" → {type, title, items:[{icon,title,description}]}
               2-column icon grid.

  "pyramid"   → {type, title, items:[{title,description}]}
               Hierarchy pyramid, first item = top/most important.

  "feature"    → {type, title, subtitle?, items:[{icon,title,description}]}
                Big title + icon features left, artwork right. Good for
                showcasing a product/service with its key capabilities.

  "statement"  → {type, title, blocks:[{heading,text}]}
                Artwork left, bold title + 2-4 content blocks right.
                Use for bold claims, positioning, or key arguments.

  "comparison" → {type, title, category?, left:{label,items:[str]}, right:{label,items:[str]}}
                Two boxes side by side: positive/pros (left) vs challenges/cons (right).
                Use for pros & cons, option A vs B, before & after.

  "grid"       → {type, title, category?, cards:[{title,description}×4]}
                2×2 card grid with left accent stripe. Use for 4 tools/features/use-cases.

  "stats"      → {type, title, category?, lead?, body?, stats:[{value,label,description}], quote?, author?}
                Big impact numbers on the right (e.g. "55%", "3x", "< 1h"),
                context text + optional quote box on the left.
                Use for KPI slides, proof points, research findings.

  "toollist"   → {type, title, category?, items:[{title,description}]}
                Art panel left, 3–4 stacked tool/option cards right, each with
                accent stripe + title + description. Use for tool comparisons,
                product overviews, or feature deep-dives.

  "roadmap"    → {type, title, category?, lead?,
                  columns:[{heading, items:[{number,title,description}]}],
                  cta?}
                Two-column numbered roadmap. Left = immediate steps,
                right = long-term. Optional bold CTA banner at the bottom.
                Use for action plans, implementation guides, timelines.
                numbers should be "01", "02", "03" etc.

Rules:
- 6–9 slides total; always start with "hero"
- Vary layout types — no two identical types in a row
- Titles: ≤ 10 words. Descriptions/bullets: 25–45 words — be specific and informative, not generic
- lead fields: 1 crisp sentence (15–25 words) that sets context for the slide content
- Available icons: star, heart, zap, shield, target, users, code, lightbulb, handshake, graduation
- Match the language of the given topic exactly
- Return ONLY valid JSON — no markdown, no code fences, no explanations

Output schema:
{"presentation_title": "...", "slides": [{...}, ...]}
"""


def call_openrouter(topic: str, model: str, api_key: str, source_text: str = None) -> dict:
    url     = "https://openrouter.ai/api/v1/chat/completions"
    headers = {
        "Authorization": f"Bearer {api_key}",
        "Content-Type":  "application/json",
        "HTTP-Referer":  "https://slidegen.local",
        "X-Title":       "SlideGen",
    }
    if source_text:
        user_msg = (
            f"Create a presentation about: {topic}\n\n"
            f"Use the following source text as the basis for the content "
            f"(summarise, structure and visualise it — do not copy verbatim):\n\n"
            f"{source_text}"
        )
    else:
        user_msg = f"Create a presentation about: {topic}"
    payload = {
        "model":    model,
        "messages": [
            {"role": "system", "content": SYSTEM_PROMPT},
            {"role": "user",   "content": user_msg},
        ],
        "temperature": 0.7,
        "max_tokens":  4000,
    }

    try:
        r = requests.post(url, headers=headers,
                          json={**payload, "response_format": {"type": "json_object"}},
                          timeout=90)
        if r.status_code not in (400, 422):
            r.raise_for_status()
            return _parse_json(r.json()["choices"][0]["message"]["content"])
    except (requests.HTTPError, json.JSONDecodeError):
        pass

    r = requests.post(url, headers=headers, json=payload, timeout=90)
    r.raise_for_status()
    return _parse_json(r.json()["choices"][0]["message"]["content"])


def _parse_json(raw: str) -> dict:
    raw = raw.strip()
    raw = re.sub(r"^```(?:json)?\s*", "", raw)
    raw = re.sub(r"\s*```$", "", raw)
    return json.loads(raw)


# ── CLI ───────────────────────────────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser(
        description="Generate AI-powered dark-themed presentations",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog=__doc__,
    )
    parser.add_argument("topic",                                                 help="Topic / idea")
    parser.add_argument("-o", "--output",    default=None,                       help="Output .pptx (auto-named if omitted)")
    parser.add_argument("-m", "--model",     default="google/gemini-flash-1.5",  help="OpenRouter model ID")
    parser.add_argument("-k", "--api-key",   default=None,                       help="OpenRouter API key (or OPENROUTER_API_KEY env var)")
    parser.add_argument("--theme",           choices=["dark", "purple"],         default="dark", help="Color theme")
    parser.add_argument("--text",            default=None,                       help="Source text file to base presentation on")
    args = parser.parse_args()

    global T
    T = THEMES[args.theme]

    api_key = args.api_key or os.environ.get("OPENROUTER_API_KEY")
    if not api_key:
        sys.exit("Error: provide --api-key or set OPENROUTER_API_KEY")

    if args.output:
        out = Path(args.output)
    else:
        safe = re.sub(r"[^\w\s-]", "", args.topic).strip()
        safe = re.sub(r"\s+", "-", safe).lower()[:50]
        out  = Path(f"{safe}.pptx")

    print(f"\n SlideGen")
    print(f" Topic  : {args.topic}")
    print(f" Model  : {args.model}")
    print(f" Theme  : {args.theme}")
    print(f" Output : {out}\n")

    source_text = None
    if args.text:
        text_path = Path(args.text)
        if not text_path.exists():
            sys.exit(f"Error: text file not found: {args.text}")
        source_text = text_path.read_text(encoding="utf-8")
        print(f" Source : {args.text} ({len(source_text)} chars)\n")

    print("Generating slide content …")
    try:
        data = call_openrouter(args.topic, args.model, api_key, source_text=source_text)
    except requests.HTTPError as e:
        sys.exit(f"API error {e.response.status_code}: {e.response.text[:300]}")
    except json.JSONDecodeError as e:
        sys.exit(f"Could not parse LLM response as JSON: {e}")

    pres_title = data.get("presentation_title", args.topic)
    slides     = data.get("slides", [])
    print(f" Title  : {pres_title}")
    print(f" Slides : {len(slides)}\n")

    print("Building slides …")
    prs = build_pptx(slides)

    prs.save(out)
    n = len(prs.slides)
    print(f"\n✓ Saved {n} slide{'s' if n != 1 else ''} → {out}")
    print("\nTo convert to PDF:")
    print(f"  macOS : open '{out}'  → File → Export as PDF")
    print(f"  Linux : libreoffice --headless --convert-to pdf '{out}'")


if __name__ == "__main__":
    main()
